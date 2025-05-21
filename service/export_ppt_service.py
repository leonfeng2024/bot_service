import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
from datetime import datetime
from typing import Dict, List, Any
import traceback

class ExportPPTService:
    def __init__(self):
        self.temp_dir = "temp"
        if not os.path.exists(self.temp_dir):
            os.makedirs(self.temp_dir)
        self.output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "output")
        # Create output directory if it doesn't exist
        os.makedirs(self.output_dir, exist_ok=True)

    def _prepare_neo4j_data(self, data):
        """
        Prepares Neo4j data for diagram generation
        
        Args:
            data: List of dictionaries or DataFrame with Neo4j relationship data
            
        Returns:
            DataFrame with standardized format for diagram generation
        """
        import pandas as pd
        
        # Convert data to DataFrame if it's a list
        if isinstance(data, list):
            df = pd.DataFrame(data)
        else:
            df = data.copy()
            
        # Check if this is Neo4j format data
        neo4j_columns = ['source_table', 'target_table', 'source_field', 'target_field']
        has_neo4j_format = any(col in df.columns for col in neo4j_columns)
        
        if has_neo4j_format:
            print("Converting Neo4j data format for diagram generation")
            # Create content field with clear relationship descriptions
            content_list = []
            
            for _, row in df.iterrows():
                source_table = row.get('source_table', '')
                target_table = row.get('target_table', '')
                source_field = row.get('source_field', '')
                target_field = row.get('target_field', '')
                
                if source_table and target_table:
                    relationship_text = f"Table {source_table} connects to table {target_table} via field {source_field} -> {target_field}"
                    content_list.append({
                        'content': relationship_text,
                        'source_table': source_table,
                        'target_table': target_table,
                        'source_field': source_field,
                        'target_field': target_field
                    })
            
            # Create new DataFrame with the content field
            if content_list:
                return pd.DataFrame(content_list)
            
        # Return original DataFrame if it doesn't match Neo4j format or conversion failed
        return df

    async def create_mermaid_diagram(self, relationships_df: pd.DataFrame, sandbox: bool = True) -> str:
        """
        Create a diagram using Mermaid from relationship data
        
        Args:
            relationships_df: DataFrame containing relationship data
            sandbox: Whether to use the sandbox mode
            
        Returns:
            Path to the generated diagram or None if no relationships
        """
        try:
            # Import os module to ensure accessibility everywhere
            import os
            import subprocess
            
            # Generate Mermaid diagram definition with improved layout settings
            mermaid_definition = """graph LR
%%{
  init: {
    'flowchart': {
      'nodeSpacing': 50,
      'rankSpacing': 80,
      'curve': 'basis',
      'nodeWidth': 200,
      'nodeHeight': 40,
      'edgeLengthFactor': '1',
      'arrowMarkerAbsolute': false,
      'htmlLabels': true
    },
    'themeVariables': {
      'fontSize': '14px',
      'fontFamily': 'Arial',
      'primaryColor': '#e6f2ff',
      'primaryTextColor': '#333',
      'primaryBorderColor': '#5d9fe4',
      'lineColor': '#4d77a5',
      'secondaryColor': '#eee',
      'tertiaryColor': '#fff',
      'arrowheadSize': '1.2'
    }
  }
}%%

classDef tableNode fill:#d7e9f7,stroke:#3c78d8,stroke-width:2px,font-size:14px,text-align:center;
classDef viewNode fill:#fff2cc,stroke:#f1c232,stroke-width:2px,font-size:14px,text-align:center;
classDef produceNode fill:#e6ffcc,stroke:#6aa84f,stroke-width:2px,font-size:14px,text-align:center;
classDef datasetNode fill:#d9d2e9,stroke:#8e7cc3,stroke-width:2px,font-size:14px,text-align:center;
classDef fieldNode fill:#fff2cc,stroke:#f1c232,stroke-width:1px,font-size:12px,text-align:left;
classDef relationshipEdge stroke:#4d77a5,stroke-width:2px;

subgraph tables["Tables"]
    direction TB
"""

            # Step 1: Extract relationship data from Neo4j format or other formats
            all_tables = set()
            all_views = set()
            all_produces = set()
            all_datasets = set()
            relationships = []
            
            # Check if the data is in Neo4j format with source_table and target_table fields
            has_neo4j_format = False
            for col in relationships_df.columns:
                if col == 'source_table' or col == 'target_table':
                    has_neo4j_format = True
                    break
            
            # Process data based on the format
            if 'content' in relationships_df.columns:
                for _, row in relationships_df.iterrows():
                    content = row.get('content', '')
                    if isinstance(content, str):
                        # Extract table names from content
                        if "Table " in content and " is related to table " in content:
                            parts = content.split(" is related to table ")
                            if len(parts) == 2:
                                source_table = parts[0].replace("Table ", "").strip()
                                
                                # Further split the target part to extract field information
                                target_parts = parts[1].split(" through field ")
                                target_table = target_parts[0].strip()
                                
                                # Add to tables set without categorization
                                all_tables.add(source_table)
                                all_tables.add(target_table)
                                
                                # Extract field information
                                field_info = ""
                                if len(target_parts) > 1:
                                    field_info = target_parts[1].strip()
                                
                                relationship = {
                                    'source': source_table,
                                    'target': target_table,
                                    'label': field_info
                                }
                                relationships.append(relationship)
            
            # If we have data in the description format
            if 'Description' in relationships_df.columns or 'description' in relationships_df.columns:
                desc_col = 'Description' if 'Description' in relationships_df.columns else 'description'
                for _, row in relationships_df.iterrows():
                    description = row.get(desc_col, '')
                    if isinstance(description, str) and " -> " in description:
                        source_field, target_field = description.split(" -> ")
                        
                        # Get source and target tables if available
                        source_table = None
                        target_table = None
                        
                        # Try to get table names from other columns
                        if 'Table' in relationships_df.columns:
                            source_table = row.get('Table')
                        elif 'source_table' in relationships_df.columns:
                            source_table = row.get('source_table')
                        
                        if 'Target Table' in relationships_df.columns:
                            target_table = row.get('Target Table')
                        elif 'target_table' in relationships_df.columns:
                            target_table = row.get('target_table')
                        
                        # If we couldn't get table names, try to extract from the fields
                        if not source_table and '_' in source_field:
                            source_table = source_field.split('_')[0]
                        
                        if not target_table and '_' in target_field:
                            target_table = target_field.split('_')[0]
                        
                        # Add tables to appropriate sets
                        if source_table:
                            all_tables.add(source_table)
                        
                        if target_table:
                            all_tables.add(target_table)
                        
                        if source_table and target_table:
                            relationship = {
                                'source': source_table,
                                'target': target_table,
                                'label': f"{source_field} -> {target_field}"
                            }
                            relationships.append(relationship)
            
            # Neo4j format with source_table and target_table fields
            if has_neo4j_format:
                print("Using Neo4j relationship data format")
                if 'source_table' in relationships_df.columns and 'target_table' in relationships_df.columns:
                    for _, row in relationships_df.iterrows():
                        source_table = row.get('source_table')
                        target_table = row.get('target_table')
                        if source_table and target_table:
                            # Add tables to appropriate sets without hardcoded categorization
                            all_tables.add(source_table)
                            all_tables.add(target_table)
                            
                            source_field = row.get('source_field', '')
                            target_field = row.get('target_field', '')
                            relationship = {
                                'source': source_table,
                                'target': target_table,
                                'label': f"{source_field} -> {target_field}"
                            }
                            relationships.append(relationship)
            
            # If no relationships found, return None instead of creating fake data
            if not relationships:
                print("No relationships found in the data. Returning None.")
                return None
            
            # Add tables to the diagram
            for table in all_tables:
                table_id = table.replace('.', '_').replace('-', '_')
                mermaid_definition += f"\n    {table_id}[<div style='padding:5px;'>{table}</div>]"
            
            mermaid_definition += "\nend\n"
            
            # Add views subgraph if we have any
            if all_views:
                mermaid_definition += "\nsubgraph views[\"Views\"]\n    direction TB\n"
                for view in all_views:
                    view_id = view.replace('.', '_').replace('-', '_')
                    mermaid_definition += f"\n    {view_id}[<div style='padding:5px;'>{view}</div>]"
                mermaid_definition += "\nend\n"
            
            # Add produce and dataset subgraph if we have any
            if all_produces or all_datasets:
                mermaid_definition += "\nsubgraph outputs[\"Produce & Dataset\"]\n    direction TB\n"
                for produce in all_produces:
                    produce_id = produce.replace('.', '_').replace('-', '_')
                    mermaid_definition += f"\n    {produce_id}[<div style='padding:5px;'>{produce}</div>]"
                
                for dataset in all_datasets:
                    dataset_id = dataset.replace('.', '_').replace('-', '_')
                    mermaid_definition += f"\n    {dataset_id}[<div style='padding:5px;'>{dataset}</div>]"
                mermaid_definition += "\nend\n"
            
            # Add relationships
            for relationship in relationships:
                source_table = relationship['source']
                target_table = relationship['target']
                label = relationship['label']
                
                # Create relationship ID
                source_id = source_table.replace('.', '_').replace('-', '_')
                target_id = target_table.replace('.', '_').replace('-', '_')
                
                # Add relationship
                mermaid_definition += f"\n{source_id} -->|{label}| {target_id}"
            
            # Set style classes for existing nodes only
            for table in all_tables:
                table_id = table.replace('.', '_').replace('-', '_')
                mermaid_definition += f"\nclass {table_id} tableNode;"
            
            for view in all_views:
                view_id = view.replace('.', '_').replace('-', '_')
                mermaid_definition += f"\nclass {view_id} viewNode;"
            
            for produce in all_produces:
                produce_id = produce.replace('.', '_').replace('-', '_')
                mermaid_definition += f"\nclass {produce_id} produceNode;"
            
            for dataset in all_datasets:
                dataset_id = dataset.replace('.', '_').replace('-', '_')
                mermaid_definition += f"\nclass {dataset_id} datasetNode;"
            
            # Ensure output directory exists
            output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'output')
            os.makedirs(output_dir, exist_ok=True)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            
            # First save mermaid definition to temporary file
            mermaid_file = os.path.join(output_dir, f"mermaid_{timestamp}.mmd")
            image_path = os.path.join(output_dir, f"relationships_{timestamp}.png")
            
            # Save mermaid definition file
            with open(mermaid_file, 'w', encoding='utf-8') as f:
                f.write(mermaid_definition)
            
            # Check for Docker vs local environment
            docker_config_path = "/app/puppeteer-config.json"
            local_config_param = ""
            
            # If we're not in Docker, don't use the config file parameter
            if not os.path.exists("/app"):
                print("Running in local environment, not using Docker puppeteer config")
                local_config_param = ""
            else:
                local_config_param = f"-p {docker_config_path}"
            
            # Use directly installed mmdc command line tool to render image - using no sandbox mode
            print("Attempting to generate diagram using mmdc command line tool...")
            try:
                # Use environment variables to ensure Puppeteer runs correctly
                env = os.environ.copy()
                env['PUPPETEER_NO_SANDBOX'] = 'true'
                env['PUPPETEER_EXECUTABLE_PATH'] = '/usr/bin/chromium' if os.path.exists('/usr/bin/chromium') else ''
                
                # Build mmdc command and execute - add no sandbox option
                cmd = ["mmdc", "-i", mermaid_file, "-o", image_path, "-b", "transparent", "-s", "3"]
                if local_config_param:
                    cmd.extend(["-p", docker_config_path])
                
                result = subprocess.run(cmd, check=True, capture_output=True, text=True, env=env)
                print(f"mmdc output: {result.stdout}")
                
                if os.path.exists(image_path) and os.path.getsize(image_path) > 1000:
                    print(f"Successfully generated diagram using mermaid-cli")
                    return image_path
                else:
                    print("Generated image file is too small or doesn't exist, trying other methods")
            except Exception as mmdc_error:
                print(f"mmdc command line error: {str(mmdc_error)}")
                if hasattr(mmdc_error, 'stderr'):
                    print(f"mmdc stderr: {mmdc_error.stderr}")
            
            # 备选方法：尝试使用NPX运行mermaid-cli，使用无沙盒模式和puppeteer配置
            print("try to use npx to run mermaid-cli...")
            try:
                # 使用npx安装并运行mmdc，添加无沙盒选项
                cmd = f"PUPPETEER_NO_SANDBOX=true npx --yes @mermaid-js/mermaid-cli mmdc -i {mermaid_file} -o {image_path} -b transparent {local_config_param} -s 3"
                result = subprocess.run(cmd, shell=True, check=True, capture_output=True, text=True, env=env)
                print(f"npx mmdc输出: {result.stdout}")
                
                if os.path.exists(image_path) and os.path.getsize(image_path) > 1000:
                    print(f"Successfully generated diagram using npx mermaid-cli")
                    return image_path
                else:
                    print("npx mmdc generated image file is too small or doesn't exist, try another method")
            except Exception as npx_error:
                print(f"npx command error: {str(npx_error)}")
                if hasattr(npx_error, 'stderr'):
                    print(f"npx stderr: {npx_error.stderr}")
            
            # 使用matplotlib作为最后的备用方法
            print("use matplotlib as a backup method to draw the relationship diagram...")
            import matplotlib.pyplot as plt
            import networkx as nx
            
            # 创建图对象
            G = nx.DiGraph()
            
            # 添加节点
            all_nodes = list(all_tables) + list(all_views) + list(all_produces) + list(all_datasets)
            for node in all_nodes:
                G.add_node(node)
            
            # 添加边
            for relationship in relationships:
                source = relationship['source']
                target = relationship['target']
                label = relationship['label']
                G.add_edge(source, target, label=label)
            
            # 图为空的情况处理 - return None instead of creating a "No relationships found" node
            if not G.nodes():
                print("no valid relationship found, returning None")
                return None
            
            # 设置图形大小
            plt.figure(figsize=(12, 9), dpi=200)
            
            # Create a custom layout with 3 columns
            pos = {}
            
            # Position tables on the left
            table_list = list(all_tables)
            for i, table in enumerate(table_list):
                pos[table] = (-3, 2 - (i * 0.5))
            
            # Position views in the middle
            view_list = list(all_views)
            for i, view in enumerate(view_list):
                pos[view] = (0, 2 - (i * 0.5))
            
            # Position produces and datasets on the right
            produce_list = list(all_produces)
            dataset_list = list(all_datasets)
            output_list = produce_list + dataset_list
            for i, output in enumerate(output_list):
                pos[output] = (3, 2 - (i * 0.5))
            
            # For any nodes not positioned yet, use spring layout
            remaining_nodes = [n for n in G.nodes() if n not in pos]
            if remaining_nodes:
                temp_pos = nx.spring_layout(G.subgraph(remaining_nodes), k=0.9, iterations=100, seed=42)
                for node, position in temp_pos.items():
                    pos[node] = position
            
            # Define node colors based on category
            node_colors = []
            for node in G.nodes():
                if node in all_tables:
                    node_colors.append('#d7e9f7')  # Blue for tables
                elif node in all_views:
                    node_colors.append('#fff2cc')  # Yellow for views
                elif node in all_produces:
                    node_colors.append('#e6ffcc')  # Green for produces
                elif node in all_datasets:
                    node_colors.append('#d9d2e9')  # Purple for datasets
                else:
                    node_colors.append('#f5f5f5')  # Gray for others
            
            # draw nodes - use larger nodes and more vivid colors
            nx.draw_networkx_nodes(G, pos, 
                                   node_size=3000, 
                                   node_color=node_colors, 
                                   edgecolors='#3c78d8', 
                                   linewidths=2)
            
            # draw edges - increase arrow size and bending degree
            nx.draw_networkx_edges(G, pos, 
                                   edge_color='#4d77a5', 
                                   width=2.0, 
                                   arrowsize=20, 
                                   arrowstyle='->', 
                                   connectionstyle='arc3,rad=0.1')
            
            # add node labels - increase font size and weight
            nx.draw_networkx_labels(G, pos, 
                                    font_size=14, 
                                    font_weight='bold',
                                    font_family='sans-serif')
            
            # add edge labels - ensure edge labels are clear and visible
            edge_labels = nx.get_edge_attributes(G, 'label')
            nx.draw_networkx_edge_labels(G, pos, 
                                         edge_labels=edge_labels, 
                                         font_size=10,
                                         font_color='#000066',
                                         bbox=dict(facecolor='white', alpha=0.7, edgecolor='none', pad=2))
            
            # save high quality image
            plt.axis('off')
            plt.tight_layout()
            plt.savefig(image_path, format='png', dpi=300, bbox_inches='tight', pad_inches=0.5, transparent=True)
            plt.close()
            
            print(f"Successfully generated diagram using matplotlib")
            print(f"Diagram saved at: {image_path}")
            
            return image_path
                
        except Exception as e:
            print(f"Error creating diagram: {str(e)}")
            traceback.print_exc()
            return None

    async def create_ppt(self, excel_file: str = "relationship.xlsx", output_file: str = "database_relationships.pptx"):
        """
        Create a PPT presentation from relationship data in Excel
        
        Args:
            excel_file: Path to the Excel file containing relationship data
            output_file: Path where the PPT file should be saved
        """
        try:
            # Read Excel file
            df = pd.read_excel(excel_file)
            print(f"Excel columns: {df.columns.tolist()}")
            
            # Prepare data for diagram generation if it's in Neo4j format
            prepared_df = self._prepare_neo4j_data(df)
            
            # Create presentation
            prs = Presentation()
            
            # Add title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "Database Relationships"
            subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            
            # Add relationship diagram slide - use a landscape slide for better fit
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Check if there are valid relationships in the data
            has_relationship_data = False
            
            # Check for Neo4j formatted data
            if 'source_table' in prepared_df.columns and 'target_table' in prepared_df.columns:
                # Check if we have actual relationship data (not empty values)
                for _, row in prepared_df.iterrows():
                    source_table = row.get('source_table')
                    target_table = row.get('target_table')
                    if source_table and target_table:
                        has_relationship_data = True
                        break
            else:
                # Check for relationship data in content field
                for _, row in prepared_df.iterrows():
                    content = row.get('content', '')
                    if content and isinstance(content, str) and ("表" in content and "关联到表" in content or 
                        "table" in content.lower() and "field" in content.lower()):
                        has_relationship_data = True
                        break
            
            # Define common slide parameters
            left = Inches(0.25)  # Reduced margin
            top = Inches(0.75)   # Reduced top margin
            width = Inches(9.5)  # Almost full width
            
            # Add title to the diagram slide first
            title_box = slide.shapes.add_textbox(left, Inches(0.1), width, Inches(0.5))
            title_frame = title_box.text_frame
            title_para = title_frame.add_paragraph()
            title_para.text = "Table Relationships Diagram"
            title_para.alignment = PP_ALIGN.CENTER
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            
            # Create and add relationship diagram only if we have relationship data
            if has_relationship_data:
                try:
                    diagram_path = await self.create_mermaid_diagram(prepared_df)
                    
                    # Check if it's a valid image file
                    if diagram_path is not None:
                        is_image = diagram_path.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))
                        
                        if is_image and os.path.exists(diagram_path) and os.path.getsize(diagram_path) > 1000:
                            slide.shapes.add_picture(diagram_path, left, top, width=width)
                        else:
                            # If image generation failed, add Japanese message
                            txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                            tf = txt_box.text_frame
                            p = tf.add_paragraph()
                            p.text = "当前検索の内容には関係図がありません"
                            p.font.size = Pt(16)
                            p.font.bold = True
                            p.alignment = PP_ALIGN.CENTER
                    else:
                        # If diagram generation completely failed, add Japanese message
                        txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                        tf = txt_box.text_frame
                        p = tf.add_paragraph()
                        p.text = "当前検索の内容には関係図がありません"
                        p.font.size = Pt(16)
                        p.font.bold = True
                        p.alignment = PP_ALIGN.CENTER
                except Exception as diagram_error:
                    print(f"Error adding diagram to PPT: {str(diagram_error)}")
                    traceback.print_exc()
                    # Add Japanese message when an error occurs
                    txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                    tf = txt_box.text_frame
                    p = tf.add_paragraph()
                    p.text = "当前検索の内容には関係図がありません"
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
            else:
                # If no relationship data, add Japanese message
                txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                tf = txt_box.text_frame
                p = tf.add_paragraph()
                p.text = "当前検索の内容には関係図がありません"
                p.font.size = Pt(16)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                print("No relationship data found, added message slide instead")
            
            # Add details slide
            table_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(table_slide_layout)
            
            # Add title to the details slide
            title_box = slide.shapes.add_textbox(Inches(0.25), Inches(0.5), Inches(9.5), Inches(0.5))
            title_frame = title_box.text_frame
            title_para = title_frame.add_paragraph()
            title_para.text = "Relationship Details"
            title_para.alignment = PP_ALIGN.CENTER
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            
            # Add table with relationship details
            rows = min(len(df) + 1, 20)  # 限制最多20行，防止表格过大
            cols = min(len(df.columns), 5)  # 限制最多5列
            top = Inches(1)
            table_width = Inches(9.5)  # Wider table
            table_height = Inches(0.3 * rows)  # Shorter rows
            
            table = slide.shapes.add_table(rows, cols, Inches(0.25), top, table_width, table_height).table
            
            # 根据实际数据列确定表头
            if 'content' in df.columns and 'description' in df.columns:
                headers = ["Content", "Description", "Created At", "Score", "Source"]
                col_map = {
                    0: 'content',
                    1: 'description',
                    2: 'created_at' if 'created_at' in df.columns else None,
                    3: 'score' if 'score' in df.columns else None,
                    4: 'source' if 'source' in df.columns else None
                }
            elif 'source_table' in df.columns and 'target_table' in df.columns:
                # Neo4j format headers
                headers = ["Source Table", "Target Table", "Source Field", "Target Field", "Description"]
                col_map = {
                    0: 'source_table',
                    1: 'target_table',
                    2: 'source_field',
                    3: 'target_field',
                    4: 'description' if 'description' in df.columns else None
                }
            else:
                # use default headers
                headers = ["Table", "First Relationship", "View", "Second Relationship", "Dataset"]
                col_map = None
            
            # set headers
            for i, header in enumerate(headers[:cols]):
                cell = table.cell(0, i)
                cell.text = header
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(11)  # Slightly smaller font
            
            # fill table data
            for i, row_data in df.iterrows():
                if i >= rows - 1:  # if exceeds the table row number limit, break the loop
                    break
                    
                for j in range(cols):
                    if col_map and col_map[j] in df.columns:
                        cell_value = str(row_data[col_map[j]]) if not pd.isna(row_data[col_map[j]]) else ""
                    else:
                        cell_value = str(row_data[j]) if j < len(df.columns) and not pd.isna(row_data[j]) else ""
                    
                    # limit cell text length, prevent cell from being too large
                    if len(cell_value) > 200:
                        cell_value = cell_value[:197] + "..."
                        
                    table.cell(i + 1, j).text = cell_value
                    # set smaller font for table data
                    paragraph = table.cell(i + 1, j).text_frame.paragraphs[0]
                    paragraph.font.size = Pt(9)  # Smaller font for data
            
            # if data rows exceed table limit, add a note
            if len(df) > rows - 1:
                note_slide = prs.slides.add_slide(table_slide_layout)
                note_box = note_slide.shapes.add_textbox(Inches(0.25), Inches(1), Inches(9.5), Inches(5))
                note_frame = note_box.text_frame
                note_para = note_frame.add_paragraph()
                note_para.text = f"注意：数据共有{len(df)}行，由于篇幅限制，只显示了前{rows-1}行。"
                note_para.font.size = Pt(14)
                
                # add title
                title_box = note_slide.shapes.add_textbox(Inches(0.25), Inches(0.1), Inches(9.5), Inches(0.5))
                title_frame = title_box.text_frame
                title_para = title_frame.add_paragraph()
                title_para.text = "数据完整性说明"
                title_para.alignment = PP_ALIGN.CENTER
                title_para.font.size = Pt(18)
                title_para.font.bold = True
            
            # Save presentation
            prs.save(output_file)
            
            # we don't delete diagram_path, keep it for reuse
            if diagram_path is not None and os.path.exists(diagram_path):
                print(f"Diagram saved at: {diagram_path}")
            print(f"PPT created at: {output_file}")
            
            return True
            
        except Exception as e:
            print(f"Error creating PPT: {str(e)}")
            traceback.print_exc()
            return False

    async def export_to_ppt(self, excel_file: str, filename_prefix: str) -> str:
        """
        Export Excel data to PowerPoint
        Returns the path to the created PowerPoint file
        """
        try:
            # Create a new presentation
            prs = Presentation()
            
            # Import pandas here to read the Excel file
            import pandas as pd
            df = pd.read_excel(excel_file)
            
            # Prepare data for diagram generation if it's in Neo4j format
            prepared_df = self._prepare_neo4j_data(df)
            
            # Add a title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "Database Query Results"
            subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
            
            # Check if there are valid relationships in the data
            has_relationship_data = False
            
            # Check for Neo4j formatted data
            if 'source_table' in prepared_df.columns and 'target_table' in prepared_df.columns:
                # Check if we have actual relationship data (not empty values)
                for _, row in prepared_df.iterrows():
                    source_table = row.get('source_table')
                    target_table = row.get('target_table')
                    if source_table and target_table:
                        has_relationship_data = True
                        break
            else:
                # Check for relationship data in content field
                for _, row in prepared_df.iterrows():
                    content = row.get('content', '')
                    if content and isinstance(content, str) and ("表" in content and "关联到表" in content or 
                        "table" in content.lower() and "field" in content.lower()):
                        has_relationship_data = True
                        break
            
            # Add diagram slide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title to the diagram slide
            left = Inches(0.25)
            top = Inches(0.75)
            width = Inches(9.5)
            
            title_box = slide.shapes.add_textbox(left, Inches(0.1), width, Inches(0.5))
            title_frame = title_box.text_frame
            title_para = title_frame.add_paragraph()
            title_para.text = "Field Relationships Diagram"
            title_para.alignment = PP_ALIGN.CENTER
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            
            # Try to create and add relationship diagram only if we have relationship data
            if has_relationship_data:
                try:
                    # Create and add relationship diagram
                    diagram_path = await self.create_mermaid_diagram(prepared_df)
                    
                    # Check if it's a valid image file
                    if diagram_path is not None:
                        is_image = diagram_path.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))
                        
                        if is_image and os.path.exists(diagram_path) and os.path.getsize(diagram_path) > 1000:
                            slide.shapes.add_picture(diagram_path, left, top, width=width)
                            print(f"Added relationship diagram to PPT")
                        else:
                            # If image generation failed, add Japanese message
                            txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                            tf = txt_box.text_frame
                            p = tf.add_paragraph()
                            p.text = "当前検索の内容には関係図がありません"
                            p.font.size = Pt(16)
                            p.font.bold = True
                            p.alignment = PP_ALIGN.CENTER
                    else:
                        # If diagram generation completely failed, add Japanese message
                        txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                        tf = txt_box.text_frame
                        p = tf.add_paragraph()
                        p.text = "当前検索の内容には関係図がありません"
                        p.font.size = Pt(16)
                        p.font.bold = True
                        p.alignment = PP_ALIGN.CENTER
                except Exception as diagram_error:
                    print(f"Error creating diagram: {str(diagram_error)}")
                    traceback.print_exc()
                    # Add Japanese message when an error occurs
                    txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                    tf = txt_box.text_frame
                    p = tf.add_paragraph()
                    p.text = "当前検索の内容には関係図がありません"
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                    print(f"Continuing with PPT creation without diagram")
            else:
                # If no relationship data, add Japanese message
                txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                tf = txt_box.text_frame
                p = tf.add_paragraph()
                p.text = "当前検索の内容には関係図がありません"
                p.font.size = Pt(16)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                print("No relationship data found, added message slide instead")
            
            # Add content slides
            # We'll create a table slide for the data
            slide_layout = prs.slide_layouts[5]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Database Relationships"
            
            # Add table
            rows, cols = df.shape
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9.0)
            height = Inches(5.0)
            
            # limit table size to avoid overflow
            rows = min(rows, 20) + 1  # +1 for header
            cols = min(cols, 10)
            
            # Add a table with headers
            shapes = slide.shapes
            table = shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set header row
            for i, col_name in enumerate(df.columns[:cols]):
                table.cell(0, i).text = str(col_name)
                # Format the header row
                for paragraph in table.cell(0, i).text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.size = Pt(12)
            
            # Fill in the data rows
            for i in range(min(rows-1, df.shape[0])):
                for j in range(min(cols, df.shape[1])):
                    cell_value = str(df.iloc[i, j]) if not pd.isna(df.iloc[i, j]) else ""
                    # limit cell text length
                    if len(cell_value) > 200:
                        cell_value = cell_value[:197] + "..."
                    table.cell(i + 1, j).text = cell_value
            
            # if data rows exceed table limit, add a note
            if df.shape[0] > rows - 1:
                note_slide = prs.slides.add_slide(slide_layout)
                note_title = note_slide.shapes.title
                note_title.text = "data integrity note"
                
                txt_box = shapes.add_textbox(left, top, width, Inches(1))
                tf = txt_box.text_frame
                p = tf.add_paragraph()
                p.text = f"caution: data has {df.shape[0]} rows, but only the first {rows-1} rows are displayed due to page limit."
                p.font.size = Pt(14)
            
            # Define output file path
            ppt_file = os.path.join(self.output_dir, f"{filename_prefix}.pptx")
            
            # Save the presentation
            prs.save(ppt_file)
            
            return ppt_file
        except Exception as e:
            print(f"Error exporting to PowerPoint: {str(e)}")
            traceback.print_exc()
            
            # try to create a basic PPT, ensure at least one output
            try:
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = "Error Creating Presentation"
                subtitle.text = f"Error: {str(e)}"
                
                # save output
                ppt_file = os.path.join(self.output_dir, f"{filename_prefix}_error.pptx")
                prs.save(ppt_file)
                
                return ppt_file
            except:
                # if even a basic PPT creation fails, raise the original exception
                raise
    
    async def append_to_ppt(self, data: List[Dict[str, Any]], ppt_file: str) -> str:
        """
        Append additional data to an existing PowerPoint file
        Returns the path to the updated PowerPoint file
        """
        try:
            # Open the existing presentation
            prs = Presentation(ppt_file)
            
            # Check if data is in Neo4j format and prepare it accordingly
            import pandas as pd
            temp_df = pd.DataFrame(data)
            prepared_df = self._prepare_neo4j_data(temp_df)
            
            # Check if data contains content field and table relationship information
            has_relationship_data = False
            is_neo4j_data = False
            is_postgresql_data = False
            
            # Check if this is PostgreSQL data that should be skipped for diagram generation
            for item in data:
                content = item.get('content', '')
                source = item.get('source', '').lower()
                if source and ('postgresql' in source or 'postgres' in source):
                    is_postgresql_data = True
                    break
                if source and 'neo4j' in source:
                    is_neo4j_data = True
            
            # Only process Neo4j data for relationship diagrams
            if is_neo4j_data and not is_postgresql_data:
                # Check for Neo4j formatted data
                if 'source_table' in prepared_df.columns and 'target_table' in prepared_df.columns:
                    # Check if we have actual relationship data (not empty values)
                    for _, row in prepared_df.iterrows():
                        source_table = row.get('source_table')
                        target_table = row.get('target_table')
                        if source_table and target_table:
                            has_relationship_data = True
                            break
                else:
                    # Check for relationship data in content field
                    for item in data:
                        content = item.get('content', '')
                        if content and isinstance(content, str) and ("表" in content and "关联到表" in content or 
                            "table" in content.lower() and "field" in content.lower()):
                            has_relationship_data = True
                            break
                
                # Add a slide for Neo4j data regardless of whether we have relationship data
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                left = Inches(0.25)
                top = Inches(0.75)
                width = Inches(9.5)
                
                # Add title to the diagram slide
                title_box = slide.shapes.add_textbox(left, Inches(0.1), width, Inches(0.5))
                title_frame = title_box.text_frame
                title_para = title_frame.add_paragraph()
                title_para.text = "Database Table Relationships"
                title_para.alignment = PP_ALIGN.CENTER
                title_para.font.size = Pt(18)
                title_para.font.bold = True
                
                if has_relationship_data:
                    try:
                        # Create and add relationship diagram
                        diagram_path = await self.create_mermaid_diagram(prepared_df)
                        
                        # Check if it's a valid image file
                        if diagram_path is not None:
                            is_image = diagram_path.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))
                            
                            if is_image and os.path.exists(diagram_path) and os.path.getsize(diagram_path) > 1000:
                                slide.shapes.add_picture(diagram_path, left, top, width=width)
                                print(f"Added new relationship diagram to PPT")
                            else:
                                # If image generation failed, add the Japanese message
                                txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                                tf = txt_box.text_frame
                                p = tf.add_paragraph()
                                p.text = "当前検索の内容には関係図がありません"
                                p.font.size = Pt(16)
                                p.font.bold = True
                                p.alignment = PP_ALIGN.CENTER
                        else:
                            # If diagram generation completely failed, add the Japanese message
                            txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                            tf = txt_box.text_frame
                            p = tf.add_paragraph()
                            p.text = "当前検索の内容には関係図がありません"
                            p.font.size = Pt(16)
                            p.font.bold = True
                            p.alignment = PP_ALIGN.CENTER
                    except Exception as diagram_error:
                        print(f"Error creating additional diagram: {str(diagram_error)}")
                        traceback.print_exc()
                        # Add the Japanese message when an error occurs
                        txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                        tf = txt_box.text_frame
                        p = tf.add_paragraph()
                        p.text = "当前検索の内容には関係図がありません"
                        p.font.size = Pt(16)
                        p.font.bold = True
                        p.alignment = PP_ALIGN.CENTER
                else:
                    # If no relationship data, add the Japanese message
                    txt_box = slide.shapes.add_textbox(left, top, width, Inches(5))
                    tf = txt_box.text_frame
                    p = tf.add_paragraph()
                    p.text = "当前検索の内容には関係図がありません"
                    p.font.size = Pt(16)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.CENTER
                    print("No relationship data found in Neo4j results, added message slide instead")
            
            # Add slides for each data item with proper content handling
            for i, item in enumerate(data):
                # Get raw content without any processing or formatting
                content = item.get('content', '')
                source = item.get('source', 'search_result')
                # Use source as title
                title_text = f"{source.capitalize()} Results"
                
                if not content:
                    continue
                
                # Add a new slide
                slide_layout = prs.slide_layouts[5]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Set slide title
                title = slide.shapes.title
                title.text = title_text
                
                # Add content text box
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9.0)
                height = Inches(5.0)
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                
                # Add original content, ensuring format preservation
                p = tf.add_paragraph()
                # Ensure content is string
                if not isinstance(content, str):
                    try:
                        import json
                        content = json.dumps(content, ensure_ascii=False, indent=2)
                    except:
                        content = str(content)
                
                p.text = content
                p.font.size = Pt(12)
            
            # Save the updated presentation
            prs.save(ppt_file)
            print(f"Successfully saved updated PPT to {ppt_file}")
            
            return ppt_file
        except Exception as e:
            print(f"Error appending to PowerPoint: {str(e)}")
            traceback.print_exc()
            # Return original file path to ensure processing continues
            return ppt_file 